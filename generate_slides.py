#!/usr/bin/env python3
"""
ぐるまっぷ 店舗オーナー様向けご案内 スライド生成スクリプト
STORE_FEATURE_GUIDE.md の内容をPPTXスライドに変換
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ブランドカラー
ORANGE = RGBColor(0xFF, 0x8C, 0x00)       # メインオレンジ
DARK_ORANGE = RGBColor(0xE0, 0x6C, 0x00)  # 濃いオレンジ
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
BLUE = RGBColor(0x21, 0x96, 0xF3)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_background(slide, color):
    """スライド背景色を設定"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    """矩形シェイプを追加"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    """角丸矩形を追加"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=BLACK,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=BLACK,
                    bullet_color=ORANGE, spacing=Pt(8)):
    """箇条書きテキストを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Meiryo"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_icon_circle(slide, left, top, size, fill_color, text="", font_size=24):
    """円形アイコンを追加"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Meiryo"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_header_bar(slide, title_text):
    """共通のヘッダーバーを追加"""
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), ORANGE)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.8),
                 title_text, font_size=32, color=WHITE, bold=True)


def add_footer(slide, page_num=None):
    """フッターを追加"""
    add_shape(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), LIGHT_GRAY)
    add_text_box(slide, Inches(0.5), Inches(7.1), Inches(5), Inches(0.4),
                 "ぐるまっぷ  店舗オーナー様向けご案内", font_size=10, color=GRAY)
    if page_num:
        add_text_box(slide, Inches(11), Inches(7.1), Inches(2), Inches(0.4),
                     str(page_num), font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)


# =====================================================
# スライド 1: タイトルスライド
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_background(slide, WHITE)

# 上部オレンジバー
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.3), ORANGE)

# メインタイトルエリア
add_shape(slide, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5), WHITE,
          border_color=ORANGE, border_width=3)

# タイトルテキスト
add_text_box(slide, Inches(2), Inches(1.8), Inches(9.333), Inches(1),
             "ぐるまっぷ", font_size=52, color=ORANGE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.0), Inches(9.333), Inches(0.8),
             "店舗オーナー様向けご案内", font_size=36, color=BLACK, bold=True,
             alignment=PP_ALIGN.CENTER)

# コンセプト
add_text_box(slide, Inches(2), Inches(4.0), Inches(9.333), Inches(0.6),
             '"知らなかった小さな名店"に出会う地図',
             font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)

# サブ情報
add_text_box(slide, Inches(2), Inches(4.8), Inches(9.333), Inches(0.6),
             "中小飲食店のための集客プラットフォームアプリ",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# 下部情報
add_text_box(slide, Inches(3), Inches(6.2), Inches(7.333), Inches(0.6),
             "初期費用 0円  |  契約縛りなし  |  iPhone / Android 対応",
             font_size=16, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# 下部オレンジバー
add_shape(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), ORANGE)


# =====================================================
# スライド 2: ぐるまっぷとは？
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "ぐるまっぷとは？")
add_footer(slide, 2)

# 3つのポイント
points = [
    ("1", "自然に発見", "マップで近くのお店を\n自然に発見してもらえます"),
    ("2", "リピーター育成", "デジタルスタンプカードで\nリピーターを育てます"),
    ("3", "効果が見える", "データ分析で\n集客の効果が見えるようになります"),
]

for i, (num, title, desc) in enumerate(points):
    x = Inches(1.0 + i * 4.0)
    y = Inches(1.8)

    # カード背景
    card = add_rounded_rect(slide, x, y, Inches(3.5), Inches(3.5), LIGHT_ORANGE)

    # 番号
    add_icon_circle(slide, x + Inches(1.25), y + Inches(0.3), Inches(1.0), ORANGE, num, 36)

    # タイトル
    add_text_box(slide, x, y + Inches(1.5), Inches(3.5), Inches(0.5),
                 title, font_size=22, color=DARK_ORANGE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # 説明
    add_text_box(slide, x + Inches(0.2), y + Inches(2.1), Inches(3.1), Inches(1.2),
                 desc, font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# 基本情報テーブル風
info_items = [
    ("対応端末", "iPhone / Android"),
    ("店舗アプリ", "専用の店舗管理アプリを無料提供"),
    ("初期費用", "0円"),
    ("契約期間", "縛りなし（いつでも解約可能）"),
]

y_base = Inches(5.6)
for i, (label, value) in enumerate(info_items):
    x = Inches(1.0 + i * 3.0)
    add_text_box(slide, x, y_base, Inches(2.8), Inches(0.3),
                 label, font_size=12, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y_base + Inches(0.3), Inches(2.8), Inches(0.4),
                 value, font_size=13, color=BLACK, alignment=PP_ALIGN.CENTER)


# =====================================================
# スライド 3: お店が得られるメリット
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "お店が得られる 3つのメリット")
add_footer(slide, 3)

merits = [
    ("1", "新しいお客さんに\n見つけてもらえる", [
        "マップ上に自動表示",
        "未訪問のお店を自動おすすめ",
        "ホームの「今日のレコメンド」で\n1ページ1店舗カード表示",
        "他の加盟店のお客さんが\nあなたのお店を発見 → 新規来店",
    ]),
    ("2", "来てくれたお客さんが\nリピーターになる", [
        "QRスキャンでスタンプが貯まる",
        "「あと少しで特典！」で再来店促進",
        "来店で自動フォロー",
        "新クーポン・投稿を\nプッシュ通知でお知らせ",
    ]),
    ("3", "集客の効果が\n数字で見える", [
        "新規顧客数・リピート率を\nリアルタイム確認",
        "月1回の担当者レポート",
        "一緒に集客を改善する\nパートナー",
        "",
    ]),
]

for i, (num, title, items) in enumerate(merits):
    x = Inches(0.5 + i * 4.2)
    y = Inches(1.5)

    # カード
    card = add_rounded_rect(slide, x, y, Inches(3.9), Inches(5.3), LIGHT_ORANGE)

    # 番号サークル
    add_icon_circle(slide, x + Inches(1.45), y + Inches(0.2), Inches(0.8), ORANGE, num, 30)

    # タイトル
    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), Inches(3.5), Inches(0.9),
                 title, font_size=18, color=DARK_ORANGE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # 箇条書き
    filtered_items = [item for item in items if item]
    for j, item in enumerate(filtered_items):
        add_text_box(slide, x + Inches(0.3), y + Inches(2.2 + j * 0.8), Inches(3.3), Inches(0.8),
                     "  " + item, font_size=13, color=BLACK)


# =====================================================
# スライド 4: スタンプカード
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "スタンプカード")
add_footer(slide, 4)

# 仕組み図
add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.4), LIGHT_ORANGE)
add_text_box(slide, Inches(1.0), Inches(1.6), Inches(11.3), Inches(0.4),
             "仕組み", font_size=14, color=ORANGE, bold=True)
add_text_box(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(0.7),
             "お客さんが来店  →  QRスキャン  →  スタンプ1個獲得  →  10個で達成  →  クーポン自動付与",
             font_size=22, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

# 詳細カード群
cards = [
    ("スタンプ数", "10個で1枚のカード達成\n（累積カウント・リセットなし）"),
    ("達成特典", "値引き型のみ\n（例:100円引き/500円引き/10%引き）\nお店が自由に金額を設定"),
    ("2枚目以降", "リセットなしで累積\n10個達成ごとに次のカードへ自動進行\n達成のたびにクーポン自動付与"),
    ("紙のカード", "不要！\nスマホで完結するので\n忘れる・なくす心配なし"),
]

for i, (title, desc) in enumerate(cards):
    x = Inches(0.8 + i * 3.1)
    y = Inches(3.3)
    card = add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.8), LIGHT_ORANGE)
    add_text_box(slide, x, y + Inches(0.2), Inches(2.8), Inches(0.4),
                 title, font_size=16, color=ORANGE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.4), Inches(2.0),
                 desc, font_size=14, color=BLACK, alignment=PP_ALIGN.CENTER)

# 物理カード移行
add_rounded_rect(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.5),
             "物理スタンプカード移行機能: 既存の紙スタンプカードのスタンプ数をアプリに引き継ぎ可能（1人1店舗あたり1回）",
             font_size=13, color=GREEN)


# =====================================================
# スライド 5: クーポン
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "クーポン機能")
add_footer(slide, 5)

# クーポン情報
coupon_items = [
    ("発行数", "1店舗あたり最大3枚まで同時発行"),
    ("種類", "割引クーポン / プレゼントクーポン / 特別オファー の3タイプ"),
    ("有効期限", "1日〜2週間の範囲で自由に設定"),
    ("発行枚数", "枚数制限あり、または無制限を選択可能"),
    ("通知", "フォロー中のお客さんにプッシュ通知でお知らせ"),
    ("効果測定", "クーポン別の利用数・推移をグラフで確認"),
]

for i, (label, value) in enumerate(coupon_items):
    y = Inches(1.5 + i * 0.65)
    add_rounded_rect(slide, Inches(0.8), y, Inches(2.5), Inches(0.55), ORANGE)
    add_text_box(slide, Inches(0.8), y + Inches(0.08), Inches(2.5), Inches(0.4),
                 label, font_size=15, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.5), y + Inches(0.08), Inches(9.0), Inches(0.4),
                 value, font_size=15, color=BLACK)

# 会計時のクーポン確認セクション
add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5), LIGHT_ORANGE)
add_text_box(slide, Inches(1.0), Inches(5.55), Inches(11.3), Inches(0.4),
             "会計時のクーポン確認（店舗スタッフ向け）",
             font_size=16, color=DARK_ORANGE, bold=True)

coupon_types = [
    ("通常クーポン", "お店が独自発行したクーポン"),
    ("コイン交換クーポン", "コイン10枚で取得した100円引き（ミッション経由）"),
    ("スタンプ達成クーポン", "スタンプ10個達成ごとに自動付与されるクーポン"),
]

for i, (ctype, cdesc) in enumerate(coupon_types):
    y = Inches(5.95 + i * 0.33)
    add_text_box(slide, Inches(1.2), y, Inches(4), Inches(0.3),
                 "  " + ctype, font_size=13, color=ORANGE, bold=True)
    add_text_box(slide, Inches(5.0), y, Inches(7.3), Inches(0.3),
                 cdesc, font_size=13, color=BLACK)


# =====================================================
# スライド 6: フォロー + プッシュ通知 & ニュース配信
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "フォロー + プッシュ通知 / ニュース配信")
add_footer(slide, 6)

# 左半分: フォロー+プッシュ通知
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0), LIGHT_ORANGE)
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.5),
             "フォロー + プッシュ通知", font_size=22, color=DARK_ORANGE, bold=True)

add_text_box(slide, Inches(0.9), Inches(2.3), Inches(5.4), Inches(0.5),
             "来店してスタンプ獲得 → 自動的にお店をフォロー",
             font_size=16, color=BLACK, bold=True)

follow_items = [
    "クーポン発行 → フォロー中のお客さんに通知",
    "投稿作成 → フォロー中のお客さんに通知",
    "通知を見て「あ、あのお店だ」→ 再来店",
]
for i, item in enumerate(follow_items):
    add_text_box(slide, Inches(1.1), Inches(3.1 + i * 0.6), Inches(5.2), Inches(0.5),
                 "  " + item, font_size=15, color=BLACK)

add_text_box(slide, Inches(0.9), Inches(5.0), Inches(5.4), Inches(0.8),
             "1回来たお客さんとの接点を\n自然に維持する仕組みです",
             font_size=15, color=GRAY, alignment=PP_ALIGN.CENTER)

# 右半分: ニュース配信
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.5),
             "ニュース配信", font_size=22, color=BLUE, bold=True)

news_items = [
    ("内容", "お知らせ・キャンペーン・季節メニューなど"),
    ("画像", "正方形（1:1）の画像を設定"),
    ("掲載期間", "開始日〜終了日を指定して自動で公開/非公開"),
    ("表示場所", "ユーザーアプリのホーム画面に\n横スクロールで表示（最大7件）"),
]
for i, (label, value) in enumerate(news_items):
    y = Inches(2.4 + i * 0.8)
    add_text_box(slide, Inches(7.2), y, Inches(2.0), Inches(0.4),
                 label, font_size=14, color=BLUE, bold=True)
    add_text_box(slide, Inches(9.2), y, Inches(3.4), Inches(0.7),
                 value, font_size=14, color=BLACK)


# =====================================================
# スライド 7: メニュー管理 & 投稿機能
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "メニュー管理 / 投稿機能")
add_footer(slide, 7)

# 左: メニュー管理
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0), LIGHT_ORANGE)
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.5),
             "メニュー管理", font_size=22, color=DARK_ORANGE, bold=True)

menu_items = [
    ("カテゴリ", "コース / 料理 / ドリンク / デザート"),
    ("表示内容", "メニュー名 / 説明 / 価格 / 写真"),
    ("オプション", "サイズ・温度など選択肢グループを作成\n追加料金も設定可能"),
    ("並び替え", "ドラッグ&ドロップで自由に順序変更"),
    ("表示先", "ユーザーアプリの店舗詳細「メニュー」タブ"),
]
for i, (label, value) in enumerate(menu_items):
    y = Inches(2.3 + i * 0.85)
    add_text_box(slide, Inches(0.9), y, Inches(2.0), Inches(0.4),
                 label, font_size=14, color=ORANGE, bold=True)
    add_text_box(slide, Inches(2.9), y, Inches(3.4), Inches(0.7),
                 value, font_size=14, color=BLACK)

# 右: 投稿機能
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.5),
             "投稿機能（プレミアムプラン）", font_size=22, color=BLUE, bold=True)

add_rounded_rect(slide, Inches(7.2), Inches(2.2), Inches(5.2), Inches(0.7), RGBColor(0xFF, 0xF9, 0xC4))
add_text_box(slide, Inches(7.4), Inches(2.3), Inches(4.8), Inches(0.5),
             "無料期間中は全機能開放のため\n投稿機能・Instagram連携もご利用いただけます",
             font_size=13, color=RGBColor(0xF5, 0x7F, 0x17), bold=True, alignment=PP_ALIGN.CENTER)

post_items = [
    ("投稿内容", "写真 + テキスト"),
    ("表示先", "ホーム画面・投稿一覧・店舗詳細"),
    ("Instagram連携", "Instagramの投稿を自動同期\n（毎日1回、時刻指定可能）"),
    ("反応", "「いいね」「コメント」を確認可能"),
]
for i, (label, value) in enumerate(post_items):
    y = Inches(3.2 + i * 0.8)
    add_text_box(slide, Inches(7.2), y, Inches(2.2), Inches(0.4),
                 label, font_size=14, color=BLUE, bold=True)
    add_text_box(slide, Inches(9.4), y, Inches(3.2), Inches(0.7),
                 value, font_size=14, color=BLACK)


# =====================================================
# スライド 8: 店舗情報の掲載
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "店舗情報の掲載")
add_footer(slide, 8)

# マップ上の表示
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.0), LIGHT_ORANGE)
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.4),
             "マップ上の表示", font_size=18, color=DARK_ORANGE, bold=True)
add_text_box(slide, Inches(0.9), Inches(2.0), Inches(5.4), Inches(0.5),
             "  お店のアイコンがマップ上にピン表示", font_size=14, color=BLACK)
add_text_box(slide, Inches(0.9), Inches(2.4), Inches(5.4), Inches(0.5),
             "  未訪問のお客さんには「おすすめ」として自動表示", font_size=14, color=BLACK)

# 店舗詳細ページ タブ
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.0), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5.6), Inches(0.4),
             "店舗詳細ページ（4タブ）", font_size=18, color=BLUE, bold=True)

tabs = [
    ("トップ", "クーポン・投稿プレビュー"),
    ("店内", "店内画像ギャラリー"),
    ("メニュー", "カテゴリ別メニュー"),
    ("投稿", "お店の最新投稿"),
]
for i, (tab_name, tab_desc) in enumerate(tabs):
    x = Inches(7.0 + i * 1.4)
    add_text_box(slide, x, Inches(2.1), Inches(1.3), Inches(0.3),
                 tab_name, font_size=13, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.4), Inches(1.3), Inches(0.5),
                 tab_desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# その他の表示情報
add_rounded_rect(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(3.0), LIGHT_GRAY)
add_text_box(slide, Inches(0.7), Inches(3.85), Inches(11.9), Inches(0.4),
             "その他の表示情報", font_size=18, color=BLACK, bold=True)

other_info = [
    "営業時間 / カテゴリ / 決済方法（30種以上に対応）",
    "臨時休業・時間変更・臨時営業の予定（今週7日間を日付付きで表示）",
    "座席情報（カウンター / テーブル / 座敷 / テラス / 個室 / ソファー）",
    "設備・サービス（駐車場 / テイクアウト / 喫煙 / Wi-Fi / バリアフリー / 子連れ / ペット）",
    "スタンプカードの進捗状況",
]
for i, info in enumerate(other_info):
    add_text_box(slide, Inches(1.0), Inches(4.4 + i * 0.45), Inches(11.5), Inches(0.4),
                 "  " + info, font_size=14, color=BLACK)


# =====================================================
# スライド 9: 分析ダッシュボード
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "分析ダッシュボード")
add_footer(slide, 9)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5),
             "店舗用アプリから、いつでもお店のデータを確認できます",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

analytics_items = [
    ("今日のサマリー", "来店者数 / 付与スタンプ数 / 新規顧客数 / クーポン利用数"),
    ("来店者推移", "日/週/月/年単位のグラフ（性別・年代でフィルター可能）"),
    ("新規顧客推移", "新規のお客さんの増減を時系列で確認"),
    ("リピート率", "週間・月間のリピート率"),
    ("クーポン効果", "クーポン別の使用者数・利用推移"),
    ("特別クーポン分析", "コイン交換クーポンの発行枚数・使用済み枚数・割引合計金額"),
    ("おすすめ表示", "レコメンドでの表示数・クリック数の推移"),
    ("送客データ", "送客元/送客先の店舗TOP5ランキング"),
]

for i, (label, value) in enumerate(analytics_items):
    row = i % 4
    col = i // 4
    x = Inches(0.5 + col * 6.4)
    y = Inches(2.0 + row * 1.2)

    card = add_rounded_rect(slide, x, y, Inches(6.1), Inches(1.0), LIGHT_ORANGE)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.05), Inches(2.0), Inches(0.4),
                 label, font_size=15, color=ORANGE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.45), Inches(5.7), Inches(0.5),
                 value, font_size=13, color=BLACK)

# 月次レポート
add_rounded_rect(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(0.7), Inches(6.05), Inches(2.5), Inches(0.4),
             "月次レポート", font_size=16, color=GREEN, bold=True)
add_text_box(slide, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.4),
             "月に1回、担当者と一緒にデータを確認 → 新規顧客数の推移 / リピート率の変化 / クーポンの効果 / 次月に向けた改善提案",
             font_size=13, color=BLACK)


# =====================================================
# スライド 10: 送客ネットワーク
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "送客ネットワーク")
add_footer(slide, 10)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.6),
             "ぐるまっぷの最大の特徴は、加盟店同士でお客さんを送り合う仕組みです。",
             font_size=20, color=BLACK, alignment=PP_ALIGN.CENTER)

# フロー図
add_rounded_rect(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5), LIGHT_ORANGE)

# A店 → B店
add_icon_circle(slide, Inches(2.0), Inches(2.8), Inches(1.2), ORANGE, "A店", 22)
add_text_box(slide, Inches(3.5), Inches(3.0), Inches(3.0), Inches(0.5),
             "→ マップでB店を発見 →", font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
add_icon_circle(slide, Inches(6.8), Inches(2.8), Inches(1.2), DARK_ORANGE, "B店", 22)
add_text_box(slide, Inches(8.3), Inches(3.0), Inches(3.0), Inches(0.5),
             "→ B店に初来店！", font_size=18, color=GREEN, bold=True)

# B店 → A店
add_icon_circle(slide, Inches(2.0), Inches(4.0), Inches(1.2), DARK_ORANGE, "B店", 22)
add_text_box(slide, Inches(3.5), Inches(4.2), Inches(3.0), Inches(0.5),
             "→ マップでA店を発見 →", font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
add_icon_circle(slide, Inches(6.8), Inches(4.0), Inches(1.2), ORANGE, "A店", 22)
add_text_box(slide, Inches(8.3), Inches(4.2), Inches(3.0), Inches(0.5),
             "→ A店に初来店！", font_size=18, color=GREEN, bold=True)

# ポイント
points = [
    "マップ上で未訪問の店舗が自動的におすすめされます",
    "あなたのお客さんが他のお店を発見し、他のお店のお客さんがあなたを発見します",
    "加盟店が増えるほど、送客の効果も大きくなります",
]
for i, point in enumerate(points):
    add_text_box(slide, Inches(1.5), Inches(5.3 + i * 0.5), Inches(10.3), Inches(0.4),
                 "  " + point, font_size=16, color=BLACK)

add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.5),
             "競合ではなく、お互いにお客さんを送り合うエコシステムです",
             font_size=18, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


# =====================================================
# スライド 11: ゲーム機能（バッジ・コイン）
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "お客さんが使うゲーム機能（お店の手間なし）")
add_footer(slide, 11)

# 左: バッジ
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.3), LIGHT_ORANGE)
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.5),
             "バッジ（全162種）", font_size=22, color=DARK_ORANGE, bold=True)
add_text_box(slide, Inches(0.9), Inches(2.2), Inches(5.4), Inches(0.5),
             "条件達成でバッジ獲得 → コレクション欲で継続利用を促進",
             font_size=15, color=GRAY)

badge_types = [
    ("スタンプ系", "スタンプ合計○個達成（15段階）"),
    ("来店系", "○回来店達成"),
    ("カテゴリ系", "カフェ/ラーメン/居酒屋 等のジャンル別\n（10グループ x 5段階）"),
    ("アクション系", "お気に入り登録 / クーポン使用 /\nコメント投稿 など"),
]
for i, (btype, bdesc) in enumerate(badge_types):
    y = Inches(2.9 + i * 0.85)
    add_text_box(slide, Inches(1.0), y, Inches(2.0), Inches(0.3),
                 btype, font_size=14, color=ORANGE, bold=True)
    add_text_box(slide, Inches(3.0), y, Inches(3.3), Inches(0.7),
                 bdesc, font_size=13, color=BLACK)

# 右: コイン
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), RGBColor(0xFF, 0xF8, 0xE1))
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.5),
             "コインシステム", font_size=22, color=RGBColor(0xF5, 0x7F, 0x17), bold=True)

coin_items = [
    ("デイリーミッション（3種）", "毎日最大3コイン"),
    ("3日連続ログイン", "+2コイン"),
    ("7日連続ログイン", "+5コイン"),
    ("30日連続ログイン", "+10コイン"),
    ("来店（スタンプ獲得時）", "+1コイン"),
    ("友達紹介", "+設定コイン数（デフォルト5）"),
]
for i, (method, coins) in enumerate(coin_items):
    y = Inches(2.3 + i * 0.55)
    add_text_box(slide, Inches(7.2), y, Inches(3.2), Inches(0.4),
                 method, font_size=12, color=BLACK)
    add_text_box(slide, Inches(10.5), y, Inches(2.1), Inches(0.4),
                 coins, font_size=12, color=RGBColor(0xF5, 0x7F, 0x17), bold=True)

add_rounded_rect(slide, Inches(7.0), Inches(5.7), Inches(5.6), Inches(0.9), RGBColor(0xFF, 0xF3, 0xE0))
add_text_box(slide, Inches(7.2), Inches(5.75), Inches(5.2), Inches(0.8),
             "使い道: 10コインで未訪問店舗の100円引きクーポンに交換\n→ まだ来たことがないお客さんの初来店きっかけに！",
             font_size=13, color=BLACK, alignment=PP_ALIGN.CENTER)


# =====================================================
# スライド 12: 対応決済方法
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "対応決済方法の表示")
add_footer(slide, 12)

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.5),
             "お店で対応している決済方法を、アプリ上でお客さんに表示できます",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

payment_categories = [
    ("現金", "現金", ORANGE),
    ("カード", "Visa / Mastercard / JCB / American Express\nDiners Club / Discover / UnionPay\nデビット / プリペイド / タッチ決済\nApple Pay / Google Pay", BLUE),
    ("電子マネー", "交通系IC / iD / QUICPay\n楽天Edy / nanaco / WAON", GREEN),
    ("QR決済", "PayPay / d払い / 楽天ペイ\nau PAY / メルペイ\nWeChat Pay / Alipay+ 他", RGBColor(0x9C, 0x27, 0xB0)),
]

for i, (category, methods, color) in enumerate(payment_categories):
    x = Inches(0.5 + i * 3.2)
    y = Inches(2.2)

    card = add_rounded_rect(slide, x, y, Inches(3.0), Inches(4.2), LIGHT_GRAY)

    # カテゴリヘッダー
    add_shape(slide, x, y, Inches(3.0), Inches(0.7), color)
    add_text_box(slide, x, y + Inches(0.1), Inches(3.0), Inches(0.5),
                 category, font_size=18, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # メソッド
    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.6), Inches(3.0),
                 methods, font_size=14, color=BLACK, alignment=PP_ALIGN.CENTER)


# =====================================================
# スライド 13: 料金プラン
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "料金プラン")
add_footer(slide, 13)

# 無料アピール
add_rounded_rect(slide, Inches(2.0), Inches(1.5), Inches(9.3), Inches(1.0), RGBColor(0xFF, 0xEB, 0xEE))
add_text_box(slide, Inches(2.2), Inches(1.55), Inches(8.9), Inches(0.8),
             "今なら先着5店舗まで無料！ 全機能開放でお使いいただけます",
             font_size=22, color=RGBColor(0xE5, 0x39, 0x35), bold=True,
             alignment=PP_ALIGN.CENTER)

# ベーシックプラン
add_rounded_rect(slide, Inches(0.8), Inches(2.8), Inches(5.8), Inches(3.5), LIGHT_ORANGE)
add_shape(slide, Inches(0.8), Inches(2.8), Inches(5.8), Inches(0.8), ORANGE)
add_text_box(slide, Inches(0.8), Inches(2.85), Inches(5.8), Inches(0.7),
             "ベーシック  月額 2,980円〜",
             font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

basic_features = [
    "マップ掲載",
    "スタンプ運用",
    "クーポン発行（3枚まで）",
    "ニュース配信",
    "メニュー管理 / 店内画像",
    "データ分析 / サポート",
]
for i, feat in enumerate(basic_features):
    add_text_box(slide, Inches(1.2), Inches(3.8 + i * 0.4), Inches(5.0), Inches(0.4),
                 "  " + feat, font_size=15, color=BLACK)

# プレミアムプラン
add_rounded_rect(slide, Inches(6.8), Inches(2.8), Inches(5.8), Inches(3.5), RGBColor(0xE3, 0xF2, 0xFD))
add_shape(slide, Inches(6.8), Inches(2.8), Inches(5.8), Inches(0.8), BLUE)
add_text_box(slide, Inches(6.8), Inches(2.85), Inches(5.8), Inches(0.7),
             "プレミアム  月額 5,980円〜",
             font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

premium_features = [
    "ベーシック全機能 +",
    "投稿機能",
    "Instagram連携",
    "",
    "※ 加盟店約30店舗到達後に導入予定",
]
for i, feat in enumerate(premium_features):
    add_text_box(slide, Inches(7.2), Inches(3.8 + i * 0.4), Inches(5.0), Inches(0.4),
                 "  " + feat, font_size=15, color=BLACK if i < 3 else GRAY)

# 1日あたり
add_rounded_rect(slide, Inches(2.5), Inches(6.5), Inches(8.3), Inches(0.7), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(2.7), Inches(6.55), Inches(7.9), Inches(0.6),
             "ベーシックプランは 1日あたり約99円 です",
             font_size=20, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# =====================================================
# スライド 14: 料金改定スケジュール
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "料金改定スケジュール")
add_footer(slide, 14)

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.5),
             "加盟店増加に伴い段階的に料金を改定する予定です（現在が最も安い料金です）",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# テーブルヘッダー
headers = ["加盟店数", "ベーシック", "プレミアム"]
for i, h in enumerate(headers):
    x = Inches(2.0 + i * 3.0)
    add_shape(slide, x, Inches(2.2), Inches(2.8), Inches(0.6), ORANGE)
    add_text_box(slide, x, Inches(2.25), Inches(2.8), Inches(0.5),
                 h, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# テーブル行
rows = [
    ("〜30店舗（現在）", "2,980円", "未提供", True),
    ("31〜60店舗", "3,980円", "5,980円", False),
    ("61〜100店舗", "4,980円", "6,980円", False),
    ("101〜150店舗", "5,980円", "8,980円", False),
    ("151店舗〜", "別途お見積もり", "別途お見積もり", False),
]

for j, (stores, basic, premium, highlight) in enumerate(rows):
    y = Inches(2.8 + j * 0.7)
    bg = LIGHT_ORANGE if highlight else WHITE
    border = ORANGE if highlight else RGBColor(0xDD, 0xDD, 0xDD)

    for i, val in enumerate([stores, basic, premium]):
        x = Inches(2.0 + i * 3.0)
        add_shape(slide, x, y, Inches(2.8), Inches(0.6), bg, border_color=border, border_width=1)
        add_text_box(slide, x, y + Inches(0.1), Inches(2.8), Inches(0.4),
                     val, font_size=15, color=DARK_ORANGE if highlight else BLACK,
                     bold=highlight, alignment=PP_ALIGN.CENTER)


# =====================================================
# スライド 15: 始め方
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "始め方（かんたん5ステップ）")
add_footer(slide, 15)

steps = [
    ("STEP 1", "アプリをダウンロード", "App Store/Google Playで\n「ぐるまっぷ 店舗」を検索", "約1分"),
    ("STEP 2", "アカウント登録", "メールアドレスとパスワードを入力\nメール認証（6桁コード）", "約5分"),
    ("STEP 3", "店舗情報を入力", "店舗名/住所/カテゴリ/営業時間\n店舗画像/設備・サービス", "10〜15分"),
    ("STEP 4", "承認を待つ", "運営チームが店舗情報を確認し\n承認します", "最短当日"),
    ("STEP 5", "利用開始！", "スタンプ達成特典/メニュー登録\n店内画像/決済方法を設定", ""),
]

for i, (step, title, desc, time_est) in enumerate(steps):
    x = Inches(0.3 + i * 2.6)
    y = Inches(1.5)

    # カード
    card_color = LIGHT_ORANGE if i < 4 else RGBColor(0xE8, 0xF5, 0xE9)
    card = add_rounded_rect(slide, x, y, Inches(2.4), Inches(4.8), card_color)

    # ステップ番号
    num_color = ORANGE if i < 4 else GREEN
    add_icon_circle(slide, x + Inches(0.6), y + Inches(0.3), Inches(1.0), num_color, str(i + 1), 30)

    # ステップ名
    add_text_box(slide, x, y + Inches(1.5), Inches(2.4), Inches(0.4),
                 step, font_size=12, color=num_color, bold=True, alignment=PP_ALIGN.CENTER)

    # タイトル
    add_text_box(slide, x, y + Inches(1.9), Inches(2.4), Inches(0.5),
                 title, font_size=16, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

    # 説明
    add_text_box(slide, x + Inches(0.1), y + Inches(2.5), Inches(2.2), Inches(1.5),
                 desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    # 所要時間
    if time_est:
        add_text_box(slide, x, y + Inches(4.2), Inches(2.4), Inches(0.4),
                     time_est, font_size=14, color=num_color, bold=True,
                     alignment=PP_ALIGN.CENTER)

# サポートメッセージ
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
             "設定はサポートが一緒にお手伝いしますので、お気軽にご相談ください",
             font_size=16, color=ORANGE, alignment=PP_ALIGN.CENTER)


# =====================================================
# スライド 16: 日常の操作
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "店舗用アプリの操作")
add_footer(slide, 16)

# 日常の操作
add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.0), LIGHT_ORANGE)
add_text_box(slide, Inches(1.0), Inches(1.55), Inches(11.3), Inches(0.4),
             "日常の操作（毎日やること）", font_size=20, color=DARK_ORANGE, bold=True)

add_text_box(slide, Inches(1.5), Inches(2.1), Inches(10.5), Inches(0.6),
             "お客さんのQRコードをスキャン → 確認 → 完了　　（約10秒）",
             font_size=22, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10.5), Inches(0.4),
             "これだけです。日常の操作はQRスキャンのみです。",
             font_size=18, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# 必要に応じてやること
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.4),
             "必要に応じてやること", font_size=18, color=BLACK, bold=True)

optional_ops = [
    ("クーポン作成", "タイトル・タイプ・有効期限・画像を入力して発行"),
    ("ニュース作成", "画像・タイトル・本文・掲載期間を入力して配信"),
    ("メニュー更新", "メニューの追加・編集・並び替え"),
    ("データ確認", "分析ダッシュボードで来店者数やリピート率を確認"),
    ("投稿作成", "写真 + テキストで投稿（プレミアム）"),
]

for i, (op, detail) in enumerate(optional_ops):
    y = Inches(4.3 + i * 0.55)
    add_rounded_rect(slide, Inches(0.8), y, Inches(2.5), Inches(0.45), ORANGE)
    add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(2.5), Inches(0.35),
                 op, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.5), y + Inches(0.05), Inches(9.0), Inches(0.35),
                 detail, font_size=14, color=BLACK)


# =====================================================
# スライド 17: サポート体制
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "サポート体制")
add_footer(slide, 17)

support_items = [
    ("ライブチャット", "アプリ内からリアルタイムで相談できます", ORANGE),
    ("メールサポート", "info@groumapapp.com\n（アプリ内フォームから送信可）", BLUE),
    ("電話サポート", "080-6050-7194\n（平日 11:00-18:00）", GREEN),
    ("月次レポート", "月に1回、データを一緒に確認して\n改善策を提案", RGBColor(0x9C, 0x27, 0xB0)),
    ("FAQ", "よくある質問をアプリ内で確認できます", RGBColor(0x79, 0x55, 0x48)),
]

for i, (title, desc, color) in enumerate(support_items):
    if i < 3:
        x = Inches(0.5 + i * 4.2)
        y = Inches(1.8)
    else:
        x = Inches(2.6 + (i - 3) * 4.2)
        y = Inches(4.5)

    card = add_rounded_rect(slide, x, y, Inches(3.8), Inches(2.2), LIGHT_GRAY)
    add_shape(slide, x, y, Inches(3.8), Inches(0.6), color)
    add_text_box(slide, x, y + Inches(0.1), Inches(3.8), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(1.2),
                 desc, font_size=15, color=BLACK, alignment=PP_ALIGN.CENTER)


# =====================================================
# スライド 18: よくあるご質問
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "よくあるご質問")
add_footer(slide, 18)

faqs = [
    ("初期費用はかかりますか？", "いいえ、初期費用は0円です。月額料金のみでご利用いただけます。"),
    ("契約期間の縛りはありますか？", "ありません。いつでも解約可能です。"),
    ("スマホが苦手でも大丈夫ですか？", "はい。日常の操作はQRコードをスキャンするだけです。初期設定はサポートが一緒にお手伝いします。"),
    ("大手の掲載サービスとの違いは？", "大手は「検索して探す」仕組み。ぐるまっぷは「マップ上で自然に見つかる」仕組みです。来店数やリピート率まで測れます。"),
    ("複数店舗を登録できますか？", "はい。1つのアカウントで複数店舗の管理が可能です。"),
    ("クーポンの費用は誰が負担？", "スタンプ達成特典・クーポンの費用はお店のご負担です。コイン交換クーポン（100円引き）の原資は現在プラットフォームが負担しています。"),
]

for i, (q, a) in enumerate(faqs):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(1.5 + row * 1.85)

    card = add_rounded_rect(slide, x, y, Inches(6.1), Inches(1.65), LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.7), Inches(0.4),
                 "Q. " + q, font_size=14, color=ORANGE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.7), Inches(1.0),
                 "A. " + a, font_size=13, color=BLACK)


# =====================================================
# スライド 19: お問い合わせ / 最終スライド
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# 上部オレンジバー
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.3), ORANGE)

# メインメッセージ
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.0),
             '"知らなかった小さな名店"として、\nあなたのお店がお客さんに発見される日を楽しみにしています。',
             font_size=28, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

# 区切り線代わり
add_shape(slide, Inches(5.0), Inches(2.8), Inches(3.3), Inches(0.05), ORANGE)

# お問い合わせ
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.6),
             "お問い合わせ", font_size=28, color=ORANGE, bold=True,
             alignment=PP_ALIGN.CENTER)

contact_items = [
    ("チャット", "アプリ内ライブチャット"),
    ("メール", "info@groumapapp.com（アプリ内お問い合わせフォームから送信可）"),
    ("電話", "080-6050-7194（平日 11:00-18:00）"),
]

for i, (method, detail) in enumerate(contact_items):
    y = Inches(4.0 + i * 0.7)
    add_rounded_rect(slide, Inches(2.5), y, Inches(8.3), Inches(0.6), LIGHT_ORANGE)
    add_text_box(slide, Inches(2.7), y + Inches(0.08), Inches(2.0), Inches(0.4),
                 method, font_size=18, color=ORANGE, bold=True)
    add_text_box(slide, Inches(5.0), y + Inches(0.08), Inches(5.6), Inches(0.4),
                 detail, font_size=16, color=BLACK)

# アピールポイント
add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.5),
             "初期費用 0円  |  契約縛りなし  |  先着5店舗 無料  |  サポート付き",
             font_size=20, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# 下部オレンジバー
add_shape(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), ORANGE)


# =====================================================
# 保存
# =====================================================
output_path = "/Users/kanekohiroki/Desktop/groumapapp/ぐるまっぷ_店舗オーナー向けご案内.pptx"
prs.save(output_path)
print(f"スライドを生成しました: {output_path}")
print(f"合計 {len(prs.slides)} スライド")
