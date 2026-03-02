#!/usr/bin/env python3
"""
ぐるまっぷ 店舗向け営業プレゼンテーション（15分版・10スライド）
STORE_SALES_PRESENTATION.md の内容をPPTXスライドに変換
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ブランドカラー
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
DARK_ORANGE = RGBColor(0xE0, 0x6C, 0x00)
DEEP_ORANGE = RGBColor(0xBF, 0x36, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x2D, 0x2D, 0x2D)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
WARM_BG = RGBColor(0xFF, 0xFA, 0xF5)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
BLUE = RGBColor(0x1E, 0x88, 0xE5)
RED = RGBColor(0xE5, 0x39, 0x35)
LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEE)
YELLOW_BG = RGBColor(0xFF, 0xF8, 0xE1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = Inches(13.333)
SH = Inches(7.5)


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, border=None, bw=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s


def rrect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT, name="Meiryo"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return tb


def multi_txt(slide, l, t, w, h, lines, sz=16, color=BLACK, spacing=Pt(6), bold=False):
    """複数行テキスト（各行をパラグラフとして追加）"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # bold指定の処理
        if isinstance(line, tuple):
            p.text = line[0]
            p.font.bold = line[1]
        else:
            p.text = line
            p.font.bold = bold
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = "Meiryo"
        p.space_after = spacing
    return tb


def circle(slide, l, t, sz, fill, text="", fsz=24):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if text:
        tf = s.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(fsz)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Meiryo"
        p.alignment = PP_ALIGN.CENTER
    return s


def slide_num(slide, num, total=10):
    txt(slide, Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4),
        f"{num}/{total}", sz=11, color=GRAY, align=PP_ALIGN.RIGHT)


def accent_bar(slide):
    rect(slide, Inches(0), Inches(7.3), SW, Inches(0.2), ORANGE)


# =====================================================
# スライド 1: 表紙（30秒）- 初頭効果
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

# 上部装飾
rect(s, Inches(0), Inches(0), SW, Inches(0.15), ORANGE)

# 中央コンテンツ
txt(s, Inches(0), Inches(1.8), SW, Inches(1.2),
    "ぐるまっぷ", sz=64, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# 区切り線
rect(s, Inches(5.0), Inches(3.1), Inches(3.333), Inches(0.06), ORANGE)

txt(s, Inches(0), Inches(3.5), SW, Inches(0.8),
    '"知らなかった小さな名店"に出会う地図', sz=30, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

txt(s, Inches(0), Inches(4.6), SW, Inches(0.6),
    "中小飲食店のための新規集客プラットフォーム", sz=22, color=GRAY, align=PP_ALIGN.CENTER)

# 下部アピール
rrect(s, Inches(3.5), Inches(5.8), Inches(6.333), Inches(0.7), LIGHT_ORANGE)
txt(s, Inches(3.5), Inches(5.88), Inches(6.333), Inches(0.5),
    "初期費用 0円  |  契約縛りなし  |  iPhone / Android 対応",
    sz=17, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

rect(s, Inches(0), Inches(7.35), SW, Inches(0.15), ORANGE)


# =====================================================
# スライド 2: 原体験+問題提起（2分）- ストーリーテリング+類似性+損失回避
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WARM_BG)

txt(s, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.7),
    "なぜ、このアプリを作ったのか", sz=36, color=DARK_ORANGE, bold=True)
rect(s, Inches(0.8), Inches(1.05), Inches(3.0), Inches(0.06), ORANGE)

# 左：ストーリー
rrect(s, Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.2), WHITE)

txt(s, Inches(0.9), Inches(1.7), Inches(7.0), Inches(0.5),
    "開発者自身が、カフェ経営者です", sz=24, color=BLACK, bold=True)

story_lines = [
    "月額約1万円の掲載サービスを契約した",
    "来店に繋がったか、測れなかった",
    "閲覧数すら少なく、費用対効果が出なかった",
]
for i, line in enumerate(story_lines):
    y = Inches(2.5 + i * 0.6)
    circle(s, Inches(1.1), y + Inches(0.05), Inches(0.35), RED if i == 1 else ORANGE, "!", 16)
    txt(s, Inches(1.7), y, Inches(6.0), Inches(0.5), line, sz=18, color=BLACK)

# 引用ボックス
rrect(s, Inches(1.0), Inches(4.5), Inches(6.5), Inches(1.8), LIGHT_RED)
rect(s, Inches(1.0), Inches(4.5), Inches(0.12), Inches(1.8), RED)
txt(s, Inches(1.4), Inches(4.7), Inches(5.8), Inches(0.5),
    "「高い掲載費を払っても、お客さんが来ない」", sz=20, color=RED, bold=True)
txt(s, Inches(1.4), Inches(5.3), Inches(5.8), Inches(0.8),
    "この悩みを解決するために、自分で作りました", sz=18, color=BLACK)

# 右：問いかけ
rrect(s, Inches(8.3), Inches(1.5), Inches(4.5), Inches(5.2), LIGHT_ORANGE)
rect(s, Inches(8.3), Inches(1.5), Inches(4.5), Inches(0.8), ORANGE)
txt(s, Inches(8.3), Inches(1.55), Inches(4.5), Inches(0.7),
    "同じ悩み、ありませんか？", sz=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

questions = [
    "掲載費に見合った\n集客ができていない",
    "来店数やリピート率が\n見えない",
    "新規のお客さんを\n増やしたい",
    "紙のスタンプカードに\n限界を感じている",
]
for i, q in enumerate(questions):
    y = Inches(2.6 + i * 1.0)
    txt(s, Inches(8.5), y, Inches(4.1), Inches(0.9),
        "  " + q, sz=15, color=BLACK)

slide_num(s, 1)
accent_bar(s)


# =====================================================
# スライド 3: 課題の可視化（1分30秒）- 恐怖喚起+コントラスト準備
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

txt(s, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.7),
    "中小飲食店の集客、どれも難しい", sz=36, color=RED, bold=True)
rect(s, Inches(0.8), Inches(1.05), Inches(3.0), Inches(0.06), RED)

# 課題テーブル
challenges = [
    ("大手掲載サービス", "月1〜2万円。大手チェーンに埋もれる", RGBColor(0xE5, 0x39, 0x35)),
    ("独自アプリ開発", "開発費100万円〜。現実的でない", RGBColor(0xD3, 0x2F, 0x2F)),
    ("紙のスタンプカード", "忘れる・なくす。続かない", RGBColor(0xC6, 0x28, 0x28)),
    ("SNS運用", "毎日の更新。本業に集中できない", RGBColor(0xB7, 0x1C, 0x1C)),
]

for i, (method, issue, color) in enumerate(challenges):
    y = Inches(1.4 + i * 1.1)

    # 方法
    rrect(s, Inches(0.8), y, Inches(3.5), Inches(0.9), LIGHT_GRAY)
    txt(s, Inches(1.0), y + Inches(0.15), Inches(3.3), Inches(0.6),
        method, sz=20, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

    # 矢印代わり
    txt(s, Inches(4.4), y + Inches(0.1), Inches(0.6), Inches(0.6),
        "→", sz=24, color=GRAY, align=PP_ALIGN.CENTER)

    # 課題
    rrect(s, Inches(5.0), y, Inches(7.8), Inches(0.9), LIGHT_RED)
    rect(s, Inches(5.0), y, Inches(0.1), Inches(0.9), color)
    txt(s, Inches(5.3), y + Inches(0.15), Inches(7.3), Inches(0.6),
        issue, sz=19, color=color, bold=True)

# このままだと...
rrect(s, Inches(1.5), Inches(5.8), Inches(10.3), Inches(1.2), RGBColor(0x2D, 0x2D, 0x2D))
txt(s, Inches(1.5), Inches(5.85), Inches(10.3), Inches(0.5),
    "このままだと...", sz=18, color=RGBColor(0xAA, 0xAA, 0xAA))
txt(s, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.5),
    "新規客が来ない → 売上が横ばい → じわじわ苦しくなる",
    sz=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

slide_num(s, 2)
accent_bar(s)


# =====================================================
# スライド 4: 解決策（2分）- コントラスト効果+チャンキング
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WARM_BG)

txt(s, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.7),
    "ぐるまっぷなら、3つの仕組みで解決", sz=36, color=DARK_ORANGE, bold=True)
rect(s, Inches(0.8), Inches(0.95), Inches(3.0), Inches(0.06), ORANGE)

solutions = [
    {
        "num": "1",
        "title": 'マップで\n"発見"される',
        "items": [
            "アプリを開くだけで、近くの未訪問\n店舗を自動おすすめ",
            "ホームの「今日のレコメンド」で\n1店舗ずつ印象に残る表示",
            "お客さんが「行ってみよう」と\n思うきっかけを作る",
        ],
        "color": ORANGE,
    },
    {
        "num": "2",
        "title": 'スタンプで\n"また来たい"を作る',
        "items": [
            "来店→QRスキャン→スタンプ獲得\nスマホで完結",
            "リセットなしで累積。10個達成ごとに\nクーポンを自動付与",
            "特典は値引き型をお店が自由に設定\n（例: 100円引き, 10%引き）",
        ],
        "color": GREEN,
    },
    {
        "num": "3",
        "title": 'データで\n"効果が見える"',
        "items": [
            "新規顧客数・リピート率・クーポン\n利用率をリアルタイム確認",
            "ミッション由来のクーポン利用枚数\n＋合計割引額も確認可能",
            "月次で一緒にデータを確認し\n改善策を考える",
        ],
        "color": BLUE,
    },
]

for i, sol in enumerate(solutions):
    x = Inches(0.3 + i * 4.3)
    y = Inches(1.3)
    c = sol["color"]

    # カード
    rrect(s, x, y, Inches(4.1), Inches(5.6), WHITE)

    # ヘッダー
    rect(s, x, y, Inches(4.1), Inches(1.6), c)
    circle(s, x + Inches(0.3), y + Inches(0.15), Inches(0.7), WHITE, sol["num"], 28)
    # 番号の色
    shapes = s.shapes
    last_shape = shapes[len(shapes) - 1]
    last_shape.text_frame.paragraphs[0].font.color.rgb = c

    txt(s, x + Inches(1.2), y + Inches(0.1), Inches(2.7), Inches(1.4),
        sol["title"], sz=22, color=WHITE, bold=True)

    # 項目
    for j, item in enumerate(sol["items"]):
        iy = y + Inches(1.9 + j * 1.15)
        txt(s, x + Inches(0.3), iy, Inches(3.5), Inches(1.0),
            item, sz=14, color=BLACK)
        if j < len(sol["items"]) - 1:
            rect(s, x + Inches(0.5), iy + Inches(1.0), Inches(3.0), Inches(0.02), LIGHT_GRAY)

# フッターメッセージ
txt(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
    '「掲載して終わり」ではなく、一緒に集客を改善するパートナーです',
    sz=17, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

slide_num(s, 3)
accent_bar(s)


# =====================================================
# スライド 5: お客さんの動き+店舗メリット（2分）- 代理体験+利得フレーミング
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

txt(s, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.7),
    "お客さんはこう動き、お店はこう得をします", sz=34, color=DARK_ORANGE, bold=True)
rect(s, Inches(0.8), Inches(0.95), Inches(3.0), Inches(0.06), ORANGE)

# フロー図（左）
flow_steps = [
    ("アプリを開く", ""),
    ("マップで未訪問店を発見", "新規客がお店を知る"),
    ("来店してQRスキャン", "来店がデータとして記録"),
    ("スタンプが貯まる", "リピートが生まれる"),
    ("10個達成→クーポン付与", "また来たい動機が強まる"),
    ("フォロー店から通知", "忘れられず、再来店"),
    ("別の店舗もおすすめ", "他店からも新規客が来る"),
]

for i, (action, benefit) in enumerate(flow_steps):
    y = Inches(1.2 + i * 0.82)

    # 左: アクション
    ac = ORANGE if i % 2 == 0 else DARK_ORANGE
    rrect(s, Inches(0.5), y, Inches(4.2), Inches(0.65), LIGHT_ORANGE)
    rect(s, Inches(0.5), y, Inches(0.1), Inches(0.65), ac)
    txt(s, Inches(0.8), y + Inches(0.1), Inches(3.7), Inches(0.45),
        action, sz=15, color=BLACK, bold=True)

    if benefit:
        # 矢印
        txt(s, Inches(4.8), y + Inches(0.05), Inches(0.8), Inches(0.5),
            "→", sz=22, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

        # 右: メリット
        rrect(s, Inches(5.6), y, Inches(3.0), Inches(0.65), LIGHT_GREEN)
        rect(s, Inches(5.6), y, Inches(0.1), Inches(0.65), GREEN)
        txt(s, Inches(5.9), y + Inches(0.1), Inches(2.5), Inches(0.45),
            benefit, sz=14, color=GREEN, bold=True)

# 右側: ゲーム機能
rrect(s, Inches(8.9), Inches(1.2), Inches(4.1), Inches(5.5), YELLOW_BG)
rect(s, Inches(8.9), Inches(1.2), Inches(4.1), Inches(0.7), RGBColor(0xF5, 0x7F, 0x17))
txt(s, Inches(8.9), Inches(1.25), Inches(4.1), Inches(0.6),
    '"つい開きたくなる"仕掛け', sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

game_items = [
    ("バッジ 162種", "コレクション欲で継続利用"),
    ("ミッション＆コイン", "10コインで100円引きクーポン\n→ 新規来店のきっかけに"),
    ("友達紹介", "紹介者・被紹介者に各+5コイン\n→ ユーザー数 = 新規客候補が増加"),
    ("自動フォロー", "来店時に自動フォロー\n→ 通知で再来店を促進"),
    ("チュートリアル", "新規登録後に自動表示\n→ 初回離脱を抑制"),
]
for i, (title, desc) in enumerate(game_items):
    gy = Inches(2.1 + i * 0.9)
    txt(s, Inches(9.1), gy, Inches(3.7), Inches(0.3),
        title, sz=13, color=RGBColor(0xF5, 0x7F, 0x17), bold=True)
    txt(s, Inches(9.1), gy + Inches(0.3), Inches(3.7), Inches(0.55),
        desc, sz=11, color=BLACK)

slide_num(s, 4)
accent_bar(s)


# =====================================================
# スライド 6: 競合比較（1分30秒）- アンカリング効果
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

txt(s, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.7),
    "大手サービスとの違い", sz=36, color=DARK_ORANGE, bold=True)
rect(s, Inches(0.8), Inches(0.95), Inches(3.0), Inches(0.06), ORANGE)

# 比較表
headers = ["比較項目", "大手掲載サービス", "ぐるまっぷ"]
header_colors = [RGBColor(0x55, 0x55, 0x55), RGBColor(0x88, 0x88, 0x88), ORANGE]

for i, (h, hc) in enumerate(zip(headers, header_colors)):
    x = Inches(0.8 + i * 4.1)
    w = Inches(3.8) if i > 0 else Inches(3.3)
    rect(s, x, Inches(1.3), w, Inches(0.7), hc)
    txt(s, x, Inches(1.35), w, Inches(0.6),
        h, sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

rows = [
    ("月額費用", "1〜5万円", "2,980円〜"),
    ("初期費用", "数万円〜", "0円"),
    ("露出", "大手に埋もれやすい", "同規模店の中で見つかる"),
    ("効果測定", "閲覧数程度", "来店数・リピート率を実測"),
    ("スタンプ", "なし / 別途契約", "標準搭載"),
    ("送客の仕組み", "検索頼み", "マップ+レコメンドで自動"),
]

for j, (item, competitor, groumap) in enumerate(rows):
    y = Inches(2.0 + j * 0.72)
    row_bg = LIGHT_GRAY if j % 2 == 0 else WHITE

    # 項目
    rect(s, Inches(0.8), y, Inches(3.3), Inches(0.65), row_bg)
    txt(s, Inches(1.0), y + Inches(0.12), Inches(3.1), Inches(0.4),
        item, sz=16, color=BLACK, bold=True)

    # 大手
    rect(s, Inches(4.9), y, Inches(3.8), Inches(0.65), row_bg)
    txt(s, Inches(5.1), y + Inches(0.12), Inches(3.4), Inches(0.4),
        competitor, sz=15, color=GRAY, align=PP_ALIGN.CENTER)

    # ぐるまっぷ
    gm_bg = LIGHT_ORANGE if j % 2 == 0 else RGBColor(0xFF, 0xF8, 0xED)
    rect(s, Inches(8.7), y, Inches(4.1), Inches(0.65), gm_bg)
    txt(s, Inches(8.9), y + Inches(0.12), Inches(3.7), Inches(0.4),
        groumap, sz=16, color=DARK_ORANGE, bold=True, align=PP_ALIGN.CENTER)

# メッセージ
rrect(s, Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.7), LIGHT_ORANGE)
txt(s, Inches(1.5), Inches(6.45), Inches(10.3), Inches(0.6),
    "今の掲載費の半額以下で、来店まで測れる仕組みが手に入ります",
    sz=20, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

slide_num(s, 5)
accent_bar(s)


# =====================================================
# スライド 7: 料金+今始めるメリット（1分30秒）- 知覚価値+希少性
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WARM_BG)

txt(s, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.7),
    "料金と、今始めるメリット", sz=36, color=DARK_ORANGE, bold=True)
rect(s, Inches(0.8), Inches(0.95), Inches(3.0), Inches(0.06), ORANGE)

# 料金カード: ベーシック
rrect(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(2.8), WHITE)
rect(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.8), ORANGE)
txt(s, Inches(0.5), Inches(1.35), Inches(6.0), Inches(0.7),
    "ベーシックプラン", sz=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(s, Inches(0.8), Inches(2.3), Inches(2.5), Inches(0.5),
    "月額", sz=16, color=GRAY)
txt(s, Inches(0.8), Inches(2.7), Inches(2.5), Inches(0.6),
    "2,980円〜", sz=32, color=ORANGE, bold=True)

txt(s, Inches(3.5), Inches(2.3), Inches(2.8), Inches(0.5),
    "1日あたり", sz=16, color=GRAY)
txt(s, Inches(3.5), Inches(2.7), Inches(2.8), Inches(0.6),
    "約99円", sz=32, color=GREEN, bold=True)

txt(s, Inches(0.8), Inches(3.4), Inches(5.5), Inches(0.5),
    "缶コーヒー1本以下の投資で集客改善", sz=14, color=GRAY)

# 含まれる機能
txt(s, Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.3),
    "マップ掲載 / スタンプ / クーポン3枚 / データ分析 / サポート",
    sz=12, color=GRAY)

# 料金カード: プレミアム
rrect(s, Inches(6.8), Inches(1.3), Inches(6.0), Inches(2.8), WHITE)
rect(s, Inches(6.8), Inches(1.3), Inches(6.0), Inches(0.8), BLUE)
txt(s, Inches(6.8), Inches(1.35), Inches(6.0), Inches(0.7),
    "プレミアムプラン", sz=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(s, Inches(7.1), Inches(2.3), Inches(2.5), Inches(0.5),
    "月額", sz=16, color=GRAY)
txt(s, Inches(7.1), Inches(2.7), Inches(2.5), Inches(0.6),
    "5,980円〜", sz=32, color=BLUE, bold=True)

txt(s, Inches(7.1), Inches(3.4), Inches(5.5), Inches(0.5),
    "ベーシック全機能 + 投稿機能 + Instagram連携", sz=14, color=GRAY)
txt(s, Inches(7.1), Inches(3.8), Inches(5.5), Inches(0.3),
    "（毎日1回の時刻指定同期 09:00〜21:00）", sz=12, color=GRAY)

# 今始めると有利な3つの理由
txt(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.5),
    "今始めると有利な 3つの理由", sz=24, color=BLACK, bold=True)

reasons = [
    ("今なら全機能無料", "先着5店舗まで無料\n投稿・Instagram連携を含む全機能が使える", RED),
    ("エリア独占", "あなたのエリアで最初なら\n近隣ユーザーを独占できる", ORANGE),
    ("料金が最安", "現在が最安値\n加盟店増加に伴い段階的に値上げ予定", GREEN),
]

for i, (title, desc, color) in enumerate(reasons):
    x = Inches(0.5 + i * 4.2)
    y = Inches(5.0)

    rrect(s, x, y, Inches(3.9), Inches(1.9), WHITE)
    rect(s, x, y, Inches(3.9), Inches(0.6), color)
    txt(s, x, y + Inches(0.08), Inches(3.9), Inches(0.45),
        title, sz=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), y + Inches(0.7), Inches(3.5), Inches(1.1),
        desc, sz=14, color=BLACK, align=PP_ALIGN.CENTER)

slide_num(s, 6)
accent_bar(s)


# =====================================================
# スライド 8: 始め方+サポート（1分）- ゼロリスクバイアス
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

txt(s, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.7),
    "始め方は簡単、サポートも万全", sz=36, color=DARK_ORANGE, bold=True)
rect(s, Inches(0.8), Inches(0.95), Inches(3.0), Inches(0.06), ORANGE)

# 3ステップ
steps = [
    ("1", "アプリをダウンロード", "App Store / Google Play\n「ぐるまっぷ 店舗」を検索", "1分"),
    ("2", "店舗情報を登録", "店舗名・住所・カテゴリ\n営業時間・画像を入力", "10〜15分"),
    ("3", "承認後、すぐ開始", "運営チームが確認し承認\nすぐに利用開始", "最短当日"),
]

for i, (num, title, desc, time_est) in enumerate(steps):
    x = Inches(0.5 + i * 4.2)
    y = Inches(1.3)

    rrect(s, x, y, Inches(3.9), Inches(2.8), LIGHT_ORANGE)
    circle(s, x + Inches(1.45), y + Inches(0.2), Inches(0.8), ORANGE, num, 30)
    txt(s, x, y + Inches(1.1), Inches(3.9), Inches(0.4),
        title, sz=18, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.3), y + Inches(1.6), Inches(3.3), Inches(0.8),
        desc, sz=14, color=GRAY, align=PP_ALIGN.CENTER)
    txt(s, x, y + Inches(2.4), Inches(3.9), Inches(0.3),
        time_est, sz=15, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# サポート体制
rrect(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(1.5), RGBColor(0xE3, 0xF2, 0xFD))
txt(s, Inches(0.7), Inches(4.45), Inches(11.9), Inches(0.4),
    "サポート体制", sz=20, color=BLUE, bold=True)

support_items = [
    "ライブチャット / メール / 電話 でいつでも相談",
    "連絡先: info@groumapapp.com / 080-6050-7194（平日 11:00-18:00）",
    "月に1回、データを一緒に確認して改善策を提案",
]
for i, item in enumerate(support_items):
    txt(s, Inches(1.0), Inches(4.95 + i * 0.35), Inches(11.5), Inches(0.35),
        "  " + item, sz=14, color=BLACK)

# お約束
rrect(s, Inches(2.0), Inches(6.2), Inches(9.3), Inches(0.9), LIGHT_GREEN)
txt(s, Inches(2.0), Inches(6.3), Inches(9.3), Inches(0.7),
    "初期費用 0円  |  契約縛りなし  |  いつでも解約可能",
    sz=24, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

slide_num(s, 7)
accent_bar(s)


# =====================================================
# スライド 9: よくある質問（2分）- 予防接種効果
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WARM_BG)

txt(s, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.7),
    "よくあるご質問", sz=36, color=DARK_ORANGE, bold=True)
rect(s, Inches(0.8), Inches(0.95), Inches(3.0), Inches(0.06), ORANGE)

faqs = [
    (
        "「うちは常連さんが多いから、新規は必要ない」",
        "常連さんが他のお店を回遊する際に、別のお店の常連さんがあなたのお店を発見します。常連さんを大切にしながら、新規も増やせます。"
    ),
    (
        "「スマホが苦手なんですが...」",
        "お店の操作はQRをスキャンするだけ。設定はサポートが一緒にやります。"
    ),
    (
        "「大手サイトとの違いは？」",
        '大手は「検索して探す」仕組み。ぐるまっぷは「自然に見つかる」仕組みです。'
    ),
    (
        "「効果が出なかったら？」",
        "契約縛りなし。月次で一緒にデータを見て改善します。リスクはゼロです。"
    ),
]

for i, (q, a) in enumerate(faqs):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 2.9)

    # カード
    rrect(s, x, y, Inches(6.0), Inches(2.6), WHITE)

    # Q
    rrect(s, x + Inches(0.2), y + Inches(0.2), Inches(5.6), Inches(0.9), LIGHT_ORANGE)
    rect(s, x + Inches(0.2), y + Inches(0.2), Inches(0.1), Inches(0.9), ORANGE)
    txt(s, x + Inches(0.5), y + Inches(0.3), Inches(0.5), Inches(0.4),
        "Q.", sz=18, color=ORANGE, bold=True)
    txt(s, x + Inches(1.0), y + Inches(0.3), Inches(4.6), Inches(0.7),
        q, sz=15, color=BLACK, bold=True)

    # A
    txt(s, Inches(0.5) + x + Inches(0.2), y + Inches(1.3), Inches(0.5), Inches(0.4),
        "A.", sz=18, color=GREEN, bold=True)
    txt(s, x + Inches(1.0), y + Inches(1.3), Inches(4.6), Inches(1.1),
        a, sz=14, color=BLACK)

slide_num(s, 8)
accent_bar(s)


# =====================================================
# スライド 10: クロージング（1分）- 選択のパラドックス回避
# =====================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)

rect(s, Inches(0), Inches(0), SW, Inches(0.15), ORANGE)

txt(s, Inches(0), Inches(0.8), SW, Inches(0.8),
    "まずは、試してみませんか？", sz=40, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

# 2つの選択肢
# A
rrect(s, Inches(0.8), Inches(2.0), Inches(5.8), Inches(2.8), LIGHT_ORANGE)
rect(s, Inches(0.8), Inches(2.0), Inches(5.8), Inches(0.8), ORANGE)
txt(s, Inches(0.8), Inches(2.05), Inches(5.8), Inches(0.7),
    "A. 今日この場で登録", sz=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1.2), Inches(3.1), Inches(5.0), Inches(1.5),
    "一緒に設定をお手伝いします\n\n15分で完了します",
    sz=18, color=BLACK, align=PP_ALIGN.CENTER)

# B
rrect(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.8), RGBColor(0xE3, 0xF2, 0xFD))
rect(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(0.8), BLUE)
txt(s, Inches(6.8), Inches(2.05), Inches(5.8), Inches(0.7),
    "B. 後日ゆっくり登録", sz=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(7.2), Inches(3.1), Inches(5.0), Inches(1.5),
    "資料をお渡しします\n\nいつでもご連絡ください",
    sz=18, color=BLACK, align=PP_ALIGN.CENTER)

# お約束
rrect(s, Inches(2.5), Inches(5.2), Inches(8.3), Inches(0.7), LIGHT_GREEN)
txt(s, Inches(2.5), Inches(5.28), Inches(8.3), Inches(0.5),
    "初期費用 0円  |  契約縛りなし  |  いつでも解約可能",
    sz=20, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# クロージングメッセージ
rect(s, Inches(0), Inches(6.1), SW, Inches(0.03), ORANGE)
txt(s, Inches(0), Inches(6.3), SW, Inches(0.8),
    '"知らなかった名店"として、\nあなたのお店が発見される日を楽しみにしています',
    sz=22, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

rect(s, Inches(0), Inches(7.35), SW, Inches(0.15), ORANGE)

slide_num(s, 9)


# =====================================================
# 保存
# =====================================================
output = "/Users/kanekohiroki/Desktop/groumapapp/ぐるまっぷ_営業プレゼン15分版.pptx"
prs.save(output)
print(f"スライドを生成しました: {output}")
print(f"合計 {len(prs.slides)} スライド")
